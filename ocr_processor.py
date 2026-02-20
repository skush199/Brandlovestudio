# ocr_processor.py

import os
import io
import json
import pdfplumber
import numpy as np
from PIL import Image
import cv2
import webcolors
from skimage.color import rgb2lab, deltaE_ciede2000
from sklearn.cluster import KMeans
import shutil


from google.cloud import vision
from google.oauth2 import service_account
from dotenv import load_dotenv

load_dotenv()

os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = os.getenv(
    "GOOGLE_APPLICATION_CREDENTIALS"
)


class GoogleVisionOCRProcessor:
    # -------------------------------------------------------
    # Detect CID garbage
    # -------------------------------------------------------
    def looks_like_cid_encoded(self, text: str) -> bool:
        return "(cid:" in (text or "").lower()

    # -------------------------------------------------------
    # Build Vision Client
    # -------------------------------------------------------
    def _build_vision_client(self, user_type: str):
        user_type = (user_type or "").strip().lower()

        if user_type == "org":
            service_account_path = os.getenv("GOOGLE_APPLICATION_CREDENTIALS")
            if not service_account_path:
                raise ValueError("GOOGLE_APPLICATION_CREDENTIALS not set")

            credentials = service_account.Credentials.from_service_account_file(
                service_account_path
            )
            return vision.ImageAnnotatorClient(credentials=credentials)

        if user_type == "byok":
            return vision.ImageAnnotatorClient()

        raise ValueError("Invalid user_type. Use 'org' or 'byok'.")

    # -------------------------------------------------------
    # OCR Image Bytes
    # -------------------------------------------------------
    def ocr_image_bytes(
        self,
        image_bytes: bytes,
        client: "vision.ImageAnnotatorClient",
        language_hints=None,
    ) -> str:
        if language_hints is None:
            language_hints = ["en"]

        image = vision.Image(content=image_bytes)
        image_context = vision.ImageContext(language_hints=language_hints)

        response = client.document_text_detection(
            image=image, image_context=image_context
        )

        if response.error.message:
            print(f"⚠️ OCR error: {response.error.message}")
            return ""

        annotation = response.full_text_annotation
        if not annotation:
            return ""

        lines = []
        if annotation.pages:
            for page in annotation.pages:
                for block in page.blocks:
                    for paragraph in block.paragraphs:
                        vertices = paragraph.bounding_box.vertices
                        y = min(v.y for v in vertices)
                        x = min(v.x for v in vertices)

                        para_text = ""
                        for word in paragraph.words:
                            for symbol in word.symbols:
                                para_text += symbol.text
                            para_text += " "

                        if para_text.strip():
                            lines.append((y, x, para_text.strip()))

        lines.sort(key=lambda l: (l[0] // 20, l[1]))

        return "\n".join([line[2] for line in lines])

    # -------------------------------------------------------
    # RGB → HEX
    # -------------------------------------------------------
    def rgb_to_hex(self, r, g, b):
        return "#{:02X}{:02X}{:02X}".format(r, g, b)

    # -------------------------------------------------------
    # Professional LAB Color Naming (Stable Version)
    # -------------------------------------------------------
    def get_color_name(self, r, g, b):
        try:
            return webcolors.rgb_to_name((r, g, b), spec="css3")
        except ValueError:
            pass

        # Convert input color to LAB
        target_rgb = np.array([[[r / 255.0, g / 255.0, b / 255.0]]])
        target_lab = rgb2lab(target_rgb)

        min_distance = float("inf")
        closest_name = None

        for name in webcolors.names("css3"):
            hex_value = webcolors.name_to_hex(name, spec="css3")
            cr, cg, cb = webcolors.hex_to_rgb(hex_value)

            comparison_rgb = np.array([[[cr / 255.0, cg / 255.0, cb / 255.0]]])
            comparison_lab = rgb2lab(comparison_rgb)

            distance = deltaE_ciede2000(target_lab, comparison_lab)[0][0]

            if distance < min_distance:
                min_distance = distance
                closest_name = name

        return closest_name

        # -------------------------------------------------------

    # WCAG Contrast Ratio
    # -------------------------------------------------------
    def _relative_luminance(self, r, g, b):
        def channel(c):
            c = c / 255.0
            return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

        R = channel(r)
        G = channel(g)
        B = channel(b)

        return 0.2126 * R + 0.7152 * G + 0.0722 * B

    def contrast_ratio(self, rgb1, rgb2):
        L1 = self._relative_luminance(*rgb1)
        L2 = self._relative_luminance(*rgb2)

        lighter = max(L1, L2)
        darker = min(L1, L2)

        return round((lighter + 0.05) / (darker + 0.05), 2)

    # -------------------------------------------------------
    # Group Similar Colors (LAB + DeltaE)
    # -------------------------------------------------------
    def group_similar_colors(self, colors, threshold=15):
        groups = []
        used = set()

        for i, c1 in enumerate(colors):
            if i in used:
                continue

            group = [c1]
            used.add(i)

            rgb1 = np.array([[[c1["r"] / 255, c1["g"] / 255, c1["b"] / 255]]])
            lab1 = rgb2lab(rgb1)

            for j, c2 in enumerate(colors):
                if j in used:
                    continue

                rgb2 = np.array([[[c2["r"] / 255, c2["g"] / 255, c2["b"] / 255]]])
                lab2 = rgb2lab(rgb2)

                delta = deltaE_ciede2000(lab1, lab2)[0][0]

                if delta < threshold:
                    group.append(c2)
                    used.add(j)

            groups.append(group)

        return groups

    # -------------------------------------------------------
    # Build Professional Palette Report
    # -------------------------------------------------------
    def build_professional_color_report(self, dominant_colors):
        groups = self.group_similar_colors(dominant_colors)

        palette = []

        for group in groups:
            total_fraction = sum(c["pixel_fraction"] for c in group)
            representative = max(group, key=lambda x: x["pixel_fraction"])

            palette.append(
                {
                    "representative_hex": representative["hex"],
                    "representative_name": representative["color_name"],
                    "grouped_colors": [c["hex"] for c in group],
                    "total_pixel_fraction": round(total_fraction, 3),
                }
            )

        palette.sort(key=lambda x: x["total_pixel_fraction"], reverse=True)
        return palette

    # -------------------------------------------------------
    # Extract Dominant Colors (NO POSITION)
    # -------------------------------------------------------
    def extract_dominant_colors(
        self, image_bytes: bytes, client: "vision.ImageAnnotatorClient"
    ):
        image = vision.Image(content=image_bytes)
        response = client.image_properties(image=image)

        colors = []

        if response.image_properties_annotation:
            for (
                color_info
            ) in response.image_properties_annotation.dominant_colors.colors:
                r = int(color_info.color.red)
                g = int(color_info.color.green)
                b = int(color_info.color.blue)

                colors.append(
                    {
                        "r": r,
                        "g": g,
                        "b": b,
                        "hex": self.rgb_to_hex(r, g, b),
                        "color_name": self.get_color_name(r, g, b),
                        "score": float(color_info.score),
                        "pixel_fraction": float(color_info.pixel_fraction),
                    }
                )

        return colors

    # -------------------------------------------------------
    # Extract ALL Colors (including white/rare colors)
    # -------------------------------------------------------
    def extract_all_colors(self, image_bytes: bytes, sample_rate: int = 4) -> list:
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        if img is None:
            return []

        img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
        sampled = img_rgb[::sample_rate, ::sample_rate].reshape(-1, 3)

        quantized = (sampled // 20) * 20
        unique, counts = np.unique(quantized, axis=0, return_counts=True)

        sorted_idx = np.argsort(-counts)
        total = len(sampled)

        all_colors = []
        for idx in sorted_idx:
            r, g, b = int(unique[idx][0]), int(unique[idx][1]), int(unique[idx][2])
            count = counts[idx]
            pct = float(count / total * 100)

            all_colors.append(
                {
                    "r": r,
                    "g": g,
                    "b": b,
                    "hex": self.rgb_to_hex(r, g, b),
                    "color_name": self.get_color_name(r, g, b),
                    "pixel_fraction": round(pct, 4),
                }
            )

        return all_colors

    # -------------------------------------------------------
    # Categorize Colors by Usage (comprehensive + readable)
    # -------------------------------------------------------
    def categorize_colors(self, image_bytes: bytes, text_boxes: list = None) -> dict:
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        if img is None:
            return {}

        height, width = img.shape[:2]
        img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)

        def get_colors_from_mask(mask, top_n=5):
            pixels = img_rgb[mask > 0]
            if len(pixels) == 0:
                return []

            quantized = (pixels // 20) * 20
            unique, counts = np.unique(quantized, axis=0, return_counts=True)
            sorted_idx = np.argsort(-counts)[:top_n]
            total = len(pixels)

            colors = []
            for idx in sorted_idx:
                r, g, b = int(unique[idx][0]), int(unique[idx][1]), int(unique[idx][2])
                pct = float(counts[idx]) / total * 100
                colors.append(
                    {
                        "hex": self.rgb_to_hex(r, g, b),
                        "color_name": self.get_color_name(r, g, b),
                        "coverage": f"{pct:.1f}",
                    }
                )
            return colors

        # Create text mask
        text_mask = np.zeros((height, width), dtype=np.uint8)
        detected_text_boxes = []

        if text_boxes:
            for x, y, w, h in text_boxes:
                cv2.rectangle(text_mask, (x, y), (x + w, y + h), 255, -1)
                detected_text_boxes.append((x, y, w, h))
        else:
            client = self._build_vision_client("org")
            image = vision.Image(content=image_bytes)
            response = client.document_text_detection(image=image)

            if response.full_text_annotation:
                for page in response.full_text_annotation.pages:
                    for block in page.blocks:
                        for paragraph in block.paragraphs:
                            vertices = paragraph.bounding_box.vertices
                            x = min(v.x for v in vertices)
                            y = min(v.y for v in vertices)
                            w = max(v.x for v in vertices) - x
                            h = max(v.y for v in vertices) - y
                            cv2.rectangle(text_mask, (x, y), (x + w, y + h), 255, -1)
                            detected_text_boxes.append((x, y, w, h))

        # 1. Background (outer edges)
        border = max(30, min(width, height) // 20)
        bg_mask = np.zeros((height, width), dtype=np.uint8)
        cv2.rectangle(
            bg_mask, (border, border), (width - border, height - border), 255, -1
        )
        bg_mask = cv2.bitwise_and(bg_mask, cv2.bitwise_not(text_mask))
        background_colors = get_colors_from_mask(bg_mask, 5)

        # 2. Text colors (inside text boxes) - get colors AND text
        text_colors = get_colors_from_mask(text_mask, 5)

        # Get text for each text color - collect OCR text
        all_text = "Sample text from image"
        try:
            client = self._build_vision_client("org")
            image = vision.Image(content=image_bytes)
            response = client.document_text_detection(image=image)

            text_by_color = {}
            text_list = []

            if response.full_text_annotation:
                # Collect all text with color matching
                for page in response.full_text_annotation.pages:
                    for block in page.blocks:
                        for paragraph in block.paragraphs:
                            para_text = ""
                            for word in paragraph.words:
                                for symbol in word.symbols:
                                    para_text += symbol.text
                                para_text += " "
                            para_text = para_text.strip()
                            if not para_text:
                                continue

                            text_list.append(para_text)

                            # Get position and color
                            vertices = paragraph.bounding_box.vertices
                            tx = min(v.x for v in vertices)
                            ty = min(v.y for v in vertices)
                            tw = max(v.x for v in vertices) - tx
                            th = max(v.y for v in vertices) - ty

                            # Sample color from text region
                            if th > 5 and tw > 5:
                                roi = img[
                                    ty : min(ty + th, height), tx : min(tx + tw, width)
                                ]
                                if roi.size > 0:
                                    avg = np.mean(roi, axis=(0, 1))
                                    r_q = (int(avg[2]) // 20) * 20
                                    g_q = (int(avg[1]) // 20) * 20
                                    b_q = (int(avg[0]) // 20) * 20
                                    hex_key = self.rgb_to_hex(r_q, g_q, b_q)
                                    if hex_key not in text_by_color:
                                        text_by_color[hex_key] = []
                                    text_by_color[hex_key].append(para_text)

                if text_list:
                    all_text = " | ".join(text_list[:5])

                # Add text to each color
                for tc in text_colors:
                    h = tc["hex"].lstrip("#")
                    r = int(h[0:2], 16)
                    g = int(h[2:4], 16)
                    b = int(h[4:6], 16)
                    r_q = (r // 20) * 20
                    g_q = (g // 20) * 20
                    b_q = (b // 20) * 20
                    hex_key = self.rgb_to_hex(r_q, g_q, b_q)

                    if hex_key in text_by_color:
                        tc["text"] = " | ".join(text_by_color[hex_key][:3])
                    elif all_text:
                        tc["text"] = all_text
        except Exception as e:
            print(f"Warning: {e}")
            for tc in text_colors:
                tc["text"] = "OCR text not available"

        # 3. UI Elements (buttons, borders, highlights)
        # Detect potential button/border regions
        edges = cv2.Canny(img, 100, 200)
        dilated = cv2.dilate(edges, np.ones((5, 5), np.uint8))
        ui_mask = cv2.bitwise_and(dilated, cv2.bitwise_not(text_mask))
        ui_colors = get_colors_from_mask(ui_mask, 5)

        # 4. Accent/Highlight areas (bright or saturated colors)
        hsv = cv2.cvtColor(img, cv2.COLOR_BGR2HSV)
        saturation = hsv[:, :, 1]
        sat_mask = (saturation > 100).astype(np.uint8) * 255
        sat_mask = cv2.bitwise_and(sat_mask, cv2.bitwise_not(text_mask))
        accent_areas = get_colors_from_mask(sat_mask, 5)

        # 5. Shadows/Dark areas
        brightness = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        dark_mask = (brightness < 80).astype(np.uint8) * 255
        dark_mask = cv2.bitwise_and(dark_mask, cv2.bitwise_not(text_mask))
        shadow_colors = get_colors_from_mask(dark_mask, 5)

        # 6. Light areas (highlights)
        light_mask = (brightness > 200).astype(np.uint8) * 255
        light_mask = cv2.bitwise_and(light_mask, cv2.bitwise_not(text_mask))
        light_colors = get_colors_from_mask(light_mask, 5)

        # Build simple output
        result = {
            "background": background_colors,
            "text": text_colors,
            "accent_foreground": accent_areas,
            "texture": ui_colors,
        }

        return result

    # -------------------------------------------------------
    # Detect Banners (large rectangular regions)
    # -------------------------------------------------------
    def detect_banners(self, image_bytes: bytes) -> list:
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        if img is None:
            return []

        height, width = img.shape[:2]
        img_gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

        # Edge detection
        edges = cv2.Canny(img_gray, 50, 150)

        # Find contours
        contours, _ = cv2.findContours(edges, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)

        banners = []

        for cnt in contours:
            x, y, w, h = cv2.boundingRect(cnt)
            aspect_ratio = float(w) / h if h > 0 else 0

            # Banner criteria: wide regions (more lenient)
            if w > width * 0.2 and h > 20 and h < height * 0.5:
                # Get dominant color in banner region
                roi = img[y : y + h, x : x + w]
                if roi.size > 0:
                    roi_rgb = cv2.cvtColor(roi, cv2.COLOR_BGR2RGB)
                    pixels = roi_rgb.reshape(-1, 3)
                    quantized = (pixels // 20) * 20
                    unique, counts = np.unique(quantized, axis=0, return_counts=True)
                    if len(unique) > 0:
                        top_idx = np.argmax(counts)
                        r, g, b = (
                            int(unique[top_idx][0]),
                            int(unique[top_idx][1]),
                            int(unique[top_idx][2]),
                        )
                        hex_code = self.rgb_to_hex(r, g, b)
                        color_name = self.get_color_name(r, g, b)

                        banners.append(
                            {
                                "x": int(x),
                                "y": int(y),
                                "w": int(w),
                                "h": int(h),
                                "hex": hex_code,
                                "color_name": color_name,
                            }
                        )

        # Remove overlapping banners (keep larger ones)
        banners = sorted(banners, key=lambda b: b["w"] * b["h"], reverse=True)
        filtered = []
        for b in banners:
            is_overlapping = False
            for f in filtered:
                if (
                    b["x"] < f["x"] + f["w"]
                    and b["x"] + b["w"] > f["x"]
                    and b["y"] < f["y"] + f["h"]
                    and b["y"] + b["h"] > f["y"]
                ):
                    is_overlapping = True
                    break
            if not is_overlapping:
                filtered.append(b)

        return filtered[:5]

        # -------------------------------------------------------
        # Find Color Positions (Bounding Boxes)
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        if img is None:
            return []

        img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
        sampled = img_rgb[::4, ::4].reshape(-1, 3)
        quantized = (sampled // 20) * 20
        unique, counts = np.unique(quantized, axis=0, return_counts=True)
        sorted_idx = np.argsort(-counts)

        return [
            {"r": int(unique[i][0]), "g": int(unique[i][1]), "b": int(unique[i][2])}
            for i in sorted_idx[:10]
        ]

    # -------------------------------------------------------
    # Find Color Positions (Bounding Boxes)
    # -------------------------------------------------------
    def find_color_positions(
        self, image_bytes: bytes, dominant_colors: list, threshold: int = 15
    ) -> list:
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        if img is None:
            return []

        height, width = img.shape[:2]
        color_positions = []

        for color in dominant_colors:
            target_color = np.array([color["b"], color["g"], color["r"]])

            lower = np.maximum(target_color - threshold, 0)
            upper = np.minimum(target_color + threshold, 255)

            mask = cv2.inRange(img, lower, upper)

            contours, _ = cv2.findContours(
                mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
            )

            regions = []
            for cnt in contours:
                if cv2.contourArea(cnt) > (width * height * 0.001):
                    x, y, w, h = cv2.boundingRect(cnt)
                    regions.append(
                        {
                            "x": int(x),
                            "y": int(y),
                            "width": int(w),
                            "height": int(h),
                            "center_x": int(x + w / 2),
                            "center_y": int(y + h / 2),
                        }
                    )

            color_positions.append(
                {
                    "hex": color["hex"],
                    "color_name": color.get("color_name", ""),
                    "regions": regions,
                    "total_regions": len(regions),
                }
            )

        return color_positions

    # -------------------------------------------------------
    # Full Image Analysis
    # -------------------------------------------------------
    def extract_full_image_analysis(
        self, image_bytes: bytes, client: "vision.ImageAnnotatorClient"
    ):
        image = vision.Image(content=image_bytes)
        result = {}

        # Get dominant colors for color_positions
        dominant_colors = self.extract_dominant_colors(
            image_bytes=image_bytes, client=client
        )

        # Get color positions (bounding boxes)
        color_pos = self.find_color_positions(
            image_bytes=image_bytes, dominant_colors=dominant_colors
        )

        # Build position lookup by hex
        pos_lookup = {c["hex"]: c.get("regions", []) for c in color_pos}

        # Add dominant colors with percentages and positions to result
        result["dominant_colors"] = [
            {
                "hex": c["hex"],
                "color_name": c.get("color_name", ""),
                "r": c["r"],
                "g": c["g"],
                "b": c["b"],
                "percentage": round(c["pixel_fraction"] * 100, 2),
                "regions": [
                    {"x": r["x"], "y": r["y"], "w": r["width"], "h": r["height"]}
                    for r in pos_lookup.get(c["hex"], [])
                    if r.get("width", 0) > 0 and r.get("height", 0) > 0
                ],
            }
            for c in dominant_colors
        ]

        # 🔥 Color categories (background, text, accent, texture)
        raw_categories = self.categorize_colors(image_bytes=image_bytes)
        result["color_categories"] = {}
        for cat_name, colors in raw_categories.items():
            result["color_categories"][cat_name] = [
                {
                    "hex": c["hex"],
                    "color_name": c.get("color_name", ""),
                    "coverage": c.get("coverage", "0%").replace("%", ""),
                }
                for c in colors
            ]

        # 🔥 Detect banners
        banners = self.detect_banners(image_bytes)

        # Always add banners key (empty if none found)
        result["banners"] = []

        if banners:
            # Get OCR text and match with banner positions
            client = self._build_vision_client("org")
            image = vision.Image(content=image_bytes)
            response = client.document_text_detection(image=image)

            # Get all text with positions from full_text_annotation
            text_data = []
            if response.full_text_annotation:
                for page in response.full_text_annotation.pages:
                    for block in page.blocks:
                        for paragraph in block.paragraphs:
                            vertices = paragraph.bounding_box.vertices
                            tx = min(v.x for v in vertices)
                            ty = min(v.y for v in vertices)
                            tw = max(v.x for v in vertices) - tx
                            th = max(v.y for v in vertices) - ty

                            # Get text from words
                            para_text = ""
                            for word in paragraph.words:
                                for symbol in word.symbols:
                                    para_text += symbol.text
                                para_text += " "

                            text_data.append(
                                {
                                    "x": tx,
                                    "y": ty,
                                    "w": tw,
                                    "h": th,
                                    "text": para_text.strip(),
                                }
                            )

            # Match text to banners
            result_banners = []
            for banner in banners:
                bx, by, bw, bh = banner["x"], banner["y"], banner["w"], banner["h"]
                banner_text = ""

                for td in text_data:
                    # Check if text is inside banner
                    if (
                        td["x"] >= bx
                        and td["y"] >= by
                        and td["x"] + td["w"] <= bx + bw
                        and td["y"] + td["h"] <= by + bh
                    ):
                        banner_text += td["text"] + " "

                result_banners.append(
                    {
                        "x": banner["x"],
                        "y": banner["y"],
                        "w": banner["w"],
                        "h": banner["h"],
                        "hex": banner["hex"],
                        "color_name": banner["color_name"],
                        "text": banner_text.strip(),
                    }
                )

            result["banners"] = result_banners

        # Extract text elements with colors (like banners but for all text)
        text_elements = []
        text_data = []

        # Use the text_data we already collected for banners
        if not text_data:
            # If banners weren't detected, get text fresh
            text_response = client.document_text_detection(image=image)
            if text_response.full_text_annotation:
                for page in text_response.full_text_annotation.pages:
                    for block in page.blocks:
                        for paragraph in block.paragraphs:
                            vertices = paragraph.bounding_box.vertices
                            tx = min(v.x for v in vertices)
                            ty = min(v.y for v in vertices)
                            tw = max(v.x for v in vertices) - tx
                            th = max(v.y for v in vertices) - ty
                            para_text = ""
                            for word in paragraph.words:
                                for symbol in word.symbols:
                                    para_text += symbol.text
                                para_text += " "
                            text_data.append(
                                {
                                    "x": tx,
                                    "y": ty,
                                    "w": tw,
                                    "h": th,
                                    "text": para_text.strip(),
                                }
                            )

        # Get image for color extraction
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        for td in text_data:
            bbox = {"x": td["x"], "y": td["y"], "width": td["w"], "height": td["h"]}
            color_info = self.extract_dominant_color_from_region(
                image_bytes, bbox, n_clusters=1
            )

            text_elements.append(
                {
                    "x": td["x"],
                    "y": td["y"],
                    "w": td["w"],
                    "h": td["h"],
                    "text": td["text"],
                    "hex": color_info["hex"],
                    "color_name": color_info["color_name"],
                    "r": color_info["r"],
                    "g": color_info["g"],
                    "b": color_info["b"],
                }
            )

        # Extract text elements with colors (following the workflow)
        text_elements = []
        text_data = []

        # Step 1: Extract ALL colors sorted by frequency (using sample_rate=4 for efficiency)
        all_image_colors = self.extract_all_colors(image_bytes, sample_rate=4)

        # Step 2: Get top N colors (e.g., top 5) - these are the most frequent in the image
        top_n_colors = all_image_colors[:5]
        top_n_hex = {c["hex"] for c in top_n_colors}

        # Step 3: Get all text regions with bounding boxes
        all_texts = self.extract_text_with_boxes(
            image_bytes=image_bytes, user_type="org"
        )

        # Step 4: For each text region, extract dominant color and check if it matches top N
        for text_item in all_texts:
            bbox = text_item["bbox"]
            color_info = self.extract_dominant_color_from_region(
                image_bytes, bbox, n_clusters=1, get_text_color=True
            )

            # Step 5: Check if this color matches any of the top N colors using LAB DeltaE
            is_top_color = False
            text_rgb = np.array(
                [
                    [
                        [
                            color_info["r"] / 255.0,
                            color_info["g"] / 255.0,
                            color_info["b"] / 255.0,
                        ]
                    ]
                ]
            )
            text_lab = rgb2lab(text_rgb)

            for top_color in top_n_colors:
                top_rgb = np.array(
                    [
                        [
                            [
                                top_color["r"] / 255.0,
                                top_color["g"] / 255.0,
                                top_color["b"] / 255.0,
                            ]
                        ]
                    ]
                )
                top_lab = rgb2lab(top_rgb)
                distance = deltaE_ciede2000(text_lab, top_lab)[0][0]
                if distance < 20:  # Threshold for color similarity
                    is_top_color = True
                    break

            # Calculate coverage percentage for the text region
            region_pixels = (
                bbox["width"] * bbox["height"]
                if bbox["width"] > 0 and bbox["height"] > 0
                else 1
            )

            text_elements.append(
                {
                    "text": text_item["text"],
                    "bounding_box": {
                        "x": bbox["x"],
                        "y": bbox["y"],
                        "w": bbox["width"],
                        "h": bbox["height"],
                    },
                    "dominant_color": {
                        "hex": color_info["hex"],
                        "color_name": color_info["color_name"],
                        "r": color_info["r"],
                        "g": color_info["g"],
                        "b": color_info["b"],
                    },
                    "priority": is_top_color,
                    "coverage_percentage": round(
                        region_pixels / (img.shape[0] * img.shape[1]) * 100, 2
                    )
                    if img is not None
                    else 0,
                }
            )

        # Step 6: Sort by priority (top color matches first), then by text length
        text_elements.sort(key=lambda x: (not x["priority"], -len(x["text"])))

        # Also add structured text that preserves visual layout
        # Sort by y (line), then by x (position in line)
        sorted_by_layout = sorted(
            text_elements,
            key=lambda x: (x["bounding_box"]["y"] // 10, x["bounding_box"]["x"]),
        )

        # Group words into lines (words on similar y positions)
        lines = []
        current_line = []
        last_y = None
        line_threshold = 15  # pixels

        for elem in sorted_by_layout:
            y = elem["bounding_box"]["y"]
            if last_y is not None and abs(y - last_y) > line_threshold:
                if current_line:
                    lines.append(current_line)
                current_line = []
            current_line.append(elem)
            last_y = y
        if current_line:
            lines.append(current_line)

        # Build structured text with layout info
        structured_text_lines = []
        for line in lines:
            line_text = " ".join([t["text"] for t in line])
            line_bbox = {
                "x": min(t["bounding_box"]["x"] for t in line),
                "y": min(t["bounding_box"]["y"] for t in line),
                "w": max(t["bounding_box"]["x"] + t["bounding_box"]["w"] for t in line)
                - min(t["bounding_box"]["x"] for t in line),
                "h": max(t["bounding_box"]["y"] + t["bounding_box"]["h"] for t in line)
                - min(t["bounding_box"]["y"] for t in line),
            }
            # Get dominant color of the line (most common or first)
            main_color = line[0]["dominant_color"]
            structured_text_lines.append(
                {
                    "text": line_text,
                    "bounding_box": line_bbox,
                    "dominant_color": main_color,
                    "word_count": len(line),
                }
            )

        # Also add the top N colors used for prioritization
        result["text_element_colors"] = [
            {
                "hex": c["hex"],
                "color_name": c["color_name"],
                "r": c["r"],
                "g": c["g"],
                "b": c["b"],
                "percentage": c["pixel_fraction"],
            }
            for c in top_n_colors
        ]

        # Add structured text that preserves visual layout (lines)
        result["structured_text"] = structured_text_lines

        # Labels
        label_response = client.label_detection(image=image)
        result["labels"] = [
            {"desc": l.description, "score": round(float(l.score), 3)}
            for l in getattr(label_response, "label_annotations", []) or []
        ]

        # Text Preview
        text_response = client.document_text_detection(image=image)
        fta = getattr(text_response, "full_text_annotation", None)

        if fta and getattr(fta, "text", None):
            result["text_preview"] = fta.text[:1500]

        # Filter out empty values
        if "banners" in result:
            result["banners"] = [
                b for b in result["banners"] if b.get("text", "").strip()
            ]

        if "color_categories" in result:
            result["color_categories"] = {
                k: v for k, v in result["color_categories"].items() if v
            }

        if "labels" in result:
            result["labels"] = [
                l for l in result["labels"] if l.get("desc", "").strip()
            ]

        return result

    # -------------------------------------------------------
    # Extract Full Page Images + Save Analysis JSON
    # -------------------------------------------------------
    def extract_images_only(
        self,
        pdf_path: str,
        output_folder: str,
        resolution: int = 500,
        user_type: str = "org",
    ):
        os.makedirs(output_folder, exist_ok=True)
        image_paths = []

        client = self._build_vision_client(user_type=user_type)

        # --- Sanity checks (avoid pdfminer "No /Root object" on non-PDF files) ---
        if not os.path.isfile(pdf_path):
            raise FileNotFoundError(f"File not found: {pdf_path}")

        with open(pdf_path, "rb") as f:
            header = f.read(5)

        if header != b"%PDF-":
            raise ValueError(f"Not a valid PDF (missing %PDF- header): {pdf_path}")

        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages):
                pdf_width = page.width
                pdf_height = page.height

                page_image = page.to_image(resolution=resolution)

                pil_image = page_image.original
                img_width, img_height = pil_image.size

                img_path = os.path.join(output_folder, f"page_{i + 1}.png")
                page_image.save(img_path)
                image_paths.append(img_path)

                print(f"🖼 Saved image: {img_path}")
                print(f"📄 PDF Size: {pdf_width} x {pdf_height}")
                print(f"🖼 Image Size: {img_width} x {img_height} before resize")

                max_size = 2000
                if img_width > max_size or img_height > max_size:
                    ratio = min(max_size / img_width, max_size / img_height)
                    new_width = int(img_width * ratio)
                    new_height = int(img_height * ratio)
                    pil_image = pil_image.resize(
                        (new_width, new_height), Image.Resampling.LANCZOS
                    )
                    img_width, img_height = pil_image.size
                    pil_image.save(img_path)
                    print(f"🖼 Resized to: {img_width} x {img_height}")

                img_bytes = io.BytesIO()
                pil_image.save(img_bytes, format="PNG")

                analysis = self.extract_full_image_analysis(
                    image_bytes=img_bytes.getvalue(), client=client
                )

                # Get text colors using new Vision API method
                text_color_result = self.extract_text_colors_background(
                    image_bytes=img_bytes.getvalue()
                )
                analysis["vision_api_text_colors"] = text_color_result["text_colors"]
                analysis["vision_api_background_colors"] = text_color_result[
                    "background_colors"
                ]
                analysis["sentences"] = text_color_result["sentences"]

                # Remove repeating/unwanted fields
                analysis.pop("text_element_colors", None)
                analysis.pop("structured_text", None)
                analysis.pop("text_preview", None)

                # 🔥 Add size metadata
                analysis["page_dimensions"] = {
                    "pdf_width_points": float(pdf_width),
                    "pdf_height_points": float(pdf_height),
                    "image_width_px": int(img_width),
                    "image_height_px": int(img_height),
                }

                json_path = img_path.replace(".png", "_analysis.json")
                with open(json_path, "w", encoding="utf-8") as f:
                    json.dump(analysis, f, indent=2, ensure_ascii=False)

                print(f"📊 Saved full analysis: {json_path}")

        return image_paths

    # <-- add at top

    def save_and_analyze_image_file(
        self, image_path: str, output_folder: str, user_type: str = "org"
    ) -> str:
        os.makedirs(output_folder, exist_ok=True)

        # copy image into output folder (this is the "store")
        stored_img_path = os.path.join(output_folder, os.path.basename(image_path))
        if os.path.abspath(image_path) != os.path.abspath(stored_img_path):
            shutil.copy2(image_path, stored_img_path)

        client = self._build_vision_client(user_type=user_type)

        with open(stored_img_path, "rb") as f:
            image_bytes = f.read()

        analysis = self.extract_full_image_analysis(
            image_bytes=image_bytes, client=client
        )

        # Get text colors using new Vision API method
        text_color_result = self.extract_text_colors_background(image_bytes=image_bytes)
        analysis["vision_api_text_colors"] = text_color_result["text_colors"]
        analysis["vision_api_background_colors"] = text_color_result[
            "background_colors"
        ]
        analysis["sentences"] = text_color_result["sentences"]

        # Remove repeating/unwanted fields
        analysis.pop("text_element_colors", None)
        analysis.pop("structured_text", None)
        analysis.pop("text_preview", None)

        json_path = os.path.splitext(stored_img_path)[0] + "_analysis.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(analysis, f, indent=2, ensure_ascii=False)

        print(f"🖼 Stored image: {stored_img_path}")
        print(f"📊 Stored analysis: {json_path}")

        return stored_img_path

    # -------------------------------------------------------
    # Extract Text From PDF
    # -------------------------------------------------------

    def extract_text_from_pdf(
        self,
        pdf_path: str,
        user_type: str = "org",
        language_hints=None,
        output_dir: str = "extracted_content",
        keep_page_breaks: bool = True,
    ) -> str:
        client = self._build_vision_client(user_type)

        # Get PDF base name (without extension)
        pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]

        # Make folder dynamic
        output_dir = os.path.join(output_dir, pdf_name)
        os.makedirs(output_dir, exist_ok=True)

        out_txt = os.path.join(output_dir, f"{pdf_name}_ocr.txt")

        all_pages_text = []

        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages):
                print(f"\n📄 Processing Page {i + 1}/{len(pdf.pages)}")

                parts = []

                normal_text = (page.extract_text() or "").strip()

                if normal_text and self.looks_like_cid_encoded(normal_text):
                    normal_text = ""

                if normal_text:
                    parts.append(normal_text)

                else:
                    page_image = page.to_image(resolution=500)
                    img_bytes = io.BytesIO()
                    page_image.save(img_bytes, format="PNG")

                    ocr_result = self.ocr_image_bytes(
                        img_bytes.getvalue(), client, language_hints=language_hints
                    )

                    if ocr_result:
                        parts.append(ocr_result)

                combined = "\n\n".join(parts).strip()

                if combined:
                    all_pages_text.append(combined)

        final_text = (
            "\n\n".join(all_pages_text) if keep_page_breaks else "".join(all_pages_text)
        )

        with open(out_txt, "w", encoding="utf-8") as f:
            f.write(final_text)

        print("✅ Text saved to:", out_txt)

        return final_text

    # should work for all formats
    # -------------------------------------------------------
    # Extract Text From Image File (Local)
    # -------------------------------------------------------
    def extract_text_from_image_file(
        self,
        image_path: str,
        user_type: str = "org",
        language_hints=None,
        output_dir: str = "extracted_content",
    ) -> str:
        if not os.path.isfile(image_path):
            raise FileNotFoundError(f"File not found: {image_path}")

        client = self._build_vision_client(user_type)

        img_name = os.path.splitext(os.path.basename(image_path))[0]
        output_dir = os.path.join(output_dir, img_name)
        os.makedirs(output_dir, exist_ok=True)

        out_txt = os.path.join(output_dir, f"{img_name}_ocr.txt")

        with open(image_path, "rb") as f:
            image_bytes = f.read()

        text = self.ocr_image_bytes(
            image_bytes=image_bytes, client=client, language_hints=language_hints
        )

        with open(out_txt, "w", encoding="utf-8") as f:
            f.write(text)

        print("✅ Text saved to:", out_txt)
        return text

    def extract_text_and_images_from_pptx(
        self,
        pptx_path: str,
        output_dir: str = "extracted_content",
        user_type: str = "org",
    ) -> dict:
        if not os.path.isfile(pptx_path):
            raise FileNotFoundError(f"File not found: {pptx_path}")

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise ImportError(
                "python-pptx is required. Install with: pip install python-pptx"
            )

        prs = Presentation(pptx_path)

        ppt_name = os.path.splitext(os.path.basename(pptx_path))[0]
        out_folder = os.path.join(output_dir, ppt_name)
        os.makedirs(out_folder, exist_ok=True)

        # build Vision client once (used for image analysis)
        client = self._build_vision_client(user_type=user_type)

        lines = []
        image_paths: list[str] = []

        for s_idx, slide in enumerate(prs.slides, start=1):
            lines.append(f"\n--- Slide {s_idx} ---")

            img_count = 0

            for shape in slide.shapes:
                # ----------------------------
                # Text boxes
                # ----------------------------
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    txt = (shape.text or "").strip()
                    if txt:
                        lines.append(txt)

                # ----------------------------
                # Tables
                # ----------------------------
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            (cell.text or "").strip() for cell in row.cells
                        ).strip()
                        if row_text:
                            lines.append(row_text)

                # ----------------------------
                # Embedded pictures
                # ----------------------------
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_count += 1
                    img = shape.image
                    ext = (img.ext or "png").lower()
                    if ext == "jpeg":
                        ext = "jpg"

                    img_path = os.path.join(
                        out_folder, f"slide_{s_idx}_img_{img_count}.{ext}"
                    )
                    with open(img_path, "wb") as f:
                        f.write(img.blob)

                    # analysis json next to image
                    analysis = self.extract_full_image_analysis(
                        image_bytes=img.blob, client=client
                    )
                    json_path = os.path.splitext(img_path)[0] + "_analysis.json"
                    with open(json_path, "w", encoding="utf-8") as jf:
                        json.dump(analysis, jf, indent=4, ensure_ascii=False)

                    image_paths.append(img_path)

                    json_path = os.path.splitext(img_path)[0] + "_analysis.json"
                    with open(json_path, "w", encoding="utf-8") as jf:
                        json.dump(analysis, jf, indent=4, ensure_ascii=False)

                    if analysis.get("text_preview"):
                        lines.append("\n--- Image OCR Text ---")
                        lines.append(analysis["text_preview"])

        extracted_text = "\n".join(lines).strip()

        out_txt = os.path.join(out_folder, f"{ppt_name}_ocr.txt")
        with open(out_txt, "w", encoding="utf-8") as f:
            f.write(extracted_text)

        print("✅ PPTX text saved to:", out_txt)

        return {
            "text": extracted_text,
            "text_path": out_txt,
            "images": image_paths,
            "output_folder": out_folder,
        }

    def extract_images_from_docx(
        self,
        docx_path: str,
        output_dir: str = "extracted_content",
        user_type: str = "org",
    ) -> list[str]:
        if not os.path.isfile(docx_path):
            raise FileNotFoundError(f"File not found: {docx_path}")

        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required. Install with: pip install python-docx"
            )

        doc_name = os.path.splitext(os.path.basename(docx_path))[0]
        out_folder = output_dir
        os.makedirs(out_folder, exist_ok=True)

        client = self._build_vision_client(user_type=user_type)

        image_paths = []
        doc = Document(docx_path)

        img_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_count += 1
                img = rel.target_part.blob
                ext = rel.target_part.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"

                img_path = os.path.join(out_folder, f"docx_img_{img_count}.{ext}")

                with open(img_path, "wb") as f:
                    f.write(img)

                analysis = self.extract_full_image_analysis(
                    image_bytes=img, client=client
                )
                json_path = os.path.splitext(img_path)[0] + "_analysis.json"
                with open(json_path, "w", encoding="utf-8") as jf:
                    json.dump(analysis, jf, indent=4, ensure_ascii=False)

                image_paths.append(img_path)
                print(f"🖼 Extracted image: {img_path}")
                print(f"📊 Saved analysis: {json_path}")

        return image_paths

    # -------------------------------------------------------
    # Extract Text with Bounding Boxes
    # -------------------------------------------------------
    def extract_text_with_boxes(
        self,
        image_path: str = None,
        image_bytes: bytes = None,
        user_type: str = "org",
    ) -> list:
        if image_path:
            with open(image_path, "rb") as f:
                image_bytes = f.read()
        elif not image_bytes:
            raise ValueError("Either image_path or image_bytes must be provided")

        client = self._build_vision_client(user_type)
        image = vision.Image(content=image_bytes)
        response = client.text_detection(image=image)

        if response.error.message:
            print(f"⚠️ OCR error: {response.error.message}")
            return []

        texts = []
        if response.text_annotations:
            for annotation in response.text_annotations[1:]:
                vertices = annotation.bounding_poly.vertices
                x = min(v.x for v in vertices)
                y = min(v.y for v in vertices)
                w = max(v.x for v in vertices) - x
                h = max(v.y for v in vertices) - y

                texts.append(
                    {
                        "text": annotation.description,
                        "bbox": {"x": x, "y": y, "width": w, "height": h},
                        "sort_y": y,
                        "sort_x": x,
                    }
                )

        texts.sort(key=lambda t: (t["sort_y"] // 10, t["sort_x"]))

        for t in texts:
            del t["sort_y"]
            del t["sort_x"]

        return texts

    # -------------------------------------------------------
    # Extract Dominant Color from Region using KMeans
    # -------------------------------------------------------
    def extract_dominant_color_from_region(
        self,
        image_bytes: bytes,
        bbox: dict,
        n_clusters: int = 1,
        get_text_color: bool = False,
    ) -> dict:
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        if img is None:
            return {"r": 0, "g": 0, "b": 0, "hex": "#000000", "color_name": "black"}

        height, width = img.shape[:2]
        x, y = max(0, bbox["x"]), max(0, bbox["y"])
        w = min(bbox["width"], width - x)
        h = min(bbox["height"], height - y)

        if w <= 0 or h <= 0:
            return {"r": 0, "g": 0, "b": 0, "hex": "#000000", "color_name": "black"}

        roi = img[y : y + h, x : x + w]

        if get_text_color:
            pixels = roi.reshape(-1, 3)

            if len(pixels) < 10:
                avg = np.mean(pixels, axis=0) if len(pixels) > 0 else [0, 0, 0]
                r, g, b = int(avg[2]), int(avg[1]), int(avg[0])
                return {
                    "r": r,
                    "g": g,
                    "b": b,
                    "hex": self.rgb_to_hex(r, g, b),
                    "color_name": self.get_color_name(r, g, b),
                }

            edge_pixels = np.vstack(
                [pixels[:10], pixels[-10:], pixels[:, :10], pixels[:, -10:]]
            )
            bg_b, bg_g, bg_r = np.median(edge_pixels, axis=0)
            bg_sum = int(bg_r) + int(bg_g) + int(bg_b)

            unique_rgb = np.unique(pixels, axis=0)
            if len(unique_rgb) == 0:
                avg = np.mean(pixels, axis=0)
                r, g, b = int(avg[2]), int(avg[1]), int(avg[0])
                return {
                    "r": r,
                    "g": g,
                    "b": b,
                    "hex": self.rgb_to_hex(r, g, b),
                    "color_name": self.get_color_name(r, g, b),
                }

            # Convert to RGB order
            unique_rgb = np.array([c[::-1].astype(int) for c in unique_rgb])

            dark_sum = int(unique_rgb.min())
            light_sum = int(unique_rgb.max())

            dark_diff = abs(dark_sum - bg_sum)
            light_diff = abs(light_sum - bg_sum)

            if light_diff > dark_diff:
                text_color = unique_rgb[unique_rgb.sum(axis=1).argmax()]
            else:
                text_color = unique_rgb[unique_rgb.sum(axis=1).argmin()]

            r, g, b = int(text_color[0]), int(text_color[1]), int(text_color[2])
            return {
                "r": int(r),
                "g": int(g),
                "b": int(b),
                "hex": self.rgb_to_hex(int(r), int(g), int(b)),
                "color_name": self.get_color_name(int(r), int(g), int(b)),
            }

        pixels = roi.reshape(-1, 3).astype(np.float32)

        if len(pixels) < 10:
            avg = np.mean(pixels, axis=0) if len(pixels) > 0 else [0, 0, 0]
            r, g, b = int(avg[2]), int(avg[1]), int(avg[0])
            return {
                "r": r,
                "g": g,
                "b": b,
                "hex": self.rgb_to_hex(r, g, b),
                "color_name": self.get_color_name(r, g, b),
            }

        criteria = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 10, 1.0)
        try:
            _, labels, centers = cv2.kmeans(
                pixels, n_clusters, None, criteria, 10, cv2.KMEANS_RANDOM_CENTERS
            )
            counts = np.bincount(labels.flatten())
            if get_text_color:
                color_idx = np.argmin(counts)
            else:
                color_idx = np.argmax(counts)
            b, g, r = centers[color_idx]
            r, g, b = int(r), int(g), int(b)
        except:
            avg = np.mean(pixels, axis=0)
            r, g, b = int(avg[2]), int(avg[1]), int(avg[0])

        return {
            "r": r,
            "g": g,
            "b": b,
            "hex": self.rgb_to_hex(r, g, b),
            "color_name": self.get_color_name(r, g, b),
        }

    # -------------------------------------------------------
    # Extract Text Colors - Main Method
    # -------------------------------------------------------
    def extract_text_colors(
        self,
        image_path: str = None,
        image_bytes: bytes = None,
        output_json_path: str = None,
        user_type: str = "org",
        visualize: bool = False,
        visualize_path: str = None,
    ) -> dict:
        if image_path:
            with open(image_path, "rb") as f:
                image_bytes = f.read()
            base_path = os.path.splitext(image_path)[0]
        elif image_bytes:
            base_path = "text_colors"
        else:
            raise ValueError("Either image_path or image_bytes must be provided")

        texts_with_boxes = self.extract_text_with_boxes(
            image_path=None, image_bytes=image_bytes, user_type=user_type
        )

        results = []
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB) if img is not None else None

        for item in texts_with_boxes:
            bbox = item["bbox"]
            color_info = self.extract_dominant_color_from_region(
                image_bytes, bbox, n_clusters=1
            )

            results.append(
                {"text": item["text"], "bounding_box": bbox, "color": color_info}
            )

        output = {"text_colors": results, "total_words": len(results)}

        if output_json_path:
            with open(output_json_path, "w", encoding="utf-8") as f:
                json.dump(output, f, indent=2, ensure_ascii=False)
            print(f"✅ Saved text colors to: {output_json_path}")

        if visualize and img is not None:
            for item in results:
                bbox = item["bounding_box"]
                color = item["color"]
                x, y, w, h = bbox["x"], bbox["y"], bbox["width"], bbox["height"]
                cv2.rectangle(
                    img, (x, y), (x + w, y + h), (color["b"], color["g"], color["r"]), 2
                )

            if not visualize_path:
                visualize_path = f"{base_path}_text_colors.png"

            cv2.imwrite(visualize_path, img)
            print(f"✅ Saved visualization to: {visualize_path}")

        return output

    # -------------------------------------------------------
    # NEW: Extract Text, Colors, Background using Vision API
    # -------------------------------------------------------
    def extract_text_colors_background(
        self,
        image_path: str = None,
        image_bytes: bytes = None,
        user_type: str = "org",
    ) -> dict:
        """
        Flow and Procedure for Text, Color, and Background Extraction:
        1. Initialize the Google Cloud Vision Client
        2. Load the Image
        3. Text Detection with Google Cloud Vision (document_text_detection)
        4. Extract Text Color and Background Color using Vision API
        5. Post-process Text Blocks into Full Sentences
        6. Return Results
        """
        # Step 1 & 2: Build client and load image
        client = self._build_vision_client(user_type)

        if image_path:
            with open(image_path, "rb") as f:
                image_bytes = f.read()
        elif not image_bytes:
            raise ValueError("Either image_path or image_bytes must be provided")

        # Step 3: Text Detection with Google Cloud Vision
        image = vision.Image(content=image_bytes)

        # Get text with full structure
        text_response = client.document_text_detection(image=image)

        # Get image properties for colors
        props_response = client.image_properties(image=image)

        full_text = ""
        if text_response.full_text_annotation:
            full_text = text_response.full_text_annotation.text

        # Step 4: Extract Text Colors and Background Colors from Vision API
        text_colors = []
        background_colors = []

        # First add white and dark colors from direct analysis
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        if img is not None:
            # Find white pixels
            white_lower = np.array([200, 200, 200])
            white_upper = np.array([255, 255, 255])
            white_mask = cv2.inRange(img, white_lower, white_upper)
            white_count = np.count_nonzero(white_mask)

            if white_count > 1000:
                text_colors.append(
                    {
                        "hex": "#FFFFFF",
                        "color_name": "white",
                        "r": 255,
                        "g": 255,
                        "b": 255,
                        "score": 0.5,
                        "pixel_fraction": white_count / (img.shape[0] * img.shape[1]),
                    }
                )

            # Find dark pixels (for dark text)
            dark_lower = np.array([0, 0, 0])
            dark_upper = np.array([60, 60, 60])
            dark_mask = cv2.inRange(img, dark_lower, dark_upper)
            dark_count = np.count_nonzero(dark_mask)

            if dark_count > 500:
                text_colors.append(
                    {
                        "hex": "#000000",
                        "color_name": "black",
                        "r": 0,
                        "g": 0,
                        "b": 0,
                        "score": 0.3,
                        "pixel_fraction": dark_count / (img.shape[0] * img.shape[1]),
                    }
                )

        if props_response.image_properties_annotation:
            colors = props_response.image_properties_annotation.dominant_colors.colors
            for color_info in colors:
                color = color_info.color
                r = int(color.red * 255) if color.red <= 1 else int(color.red)
                g = int(color.green * 255) if color.green <= 1 else int(color.green)
                b = int(color.blue * 255) if color.blue <= 1 else int(color.blue)
                rgb = {"r": r, "g": g, "b": b}
                hex_code = self.rgb_to_hex(rgb["r"], rgb["g"], rgb["b"])
                text_colors.append(
                    {
                        "hex": hex_code,
                        "color_name": self.get_color_name(rgb["r"], rgb["g"], rgb["b"]),
                        "r": rgb["r"],
                        "g": rgb["g"],
                        "b": rgb["b"],
                        "score": color_info.score,
                        "pixel_fraction": color_info.pixel_fraction,
                    }
                )

        # Also get background as most common color
        if text_colors:
            background_colors = text_colors[:5]  # Top 5 colors as potential backgrounds

        # Step 5: Post-process Text Blocks into Full Sentences
        sentences = []

        if text_response.full_text_annotation:
            for page in text_response.full_text_annotation.pages:
                for block in page.blocks:
                    for paragraph in block.paragraphs:
                        vertices = paragraph.bounding_box.vertices
                        x = min(v.x for v in vertices)
                        y = min(v.y for v in vertices)
                        w = max(v.x for v in vertices) - x
                        h = max(v.y for v in vertices) - y

                        # Get text from words
                        para_text = ""
                        for word in paragraph.words:
                            for symbol in word.symbols:
                                para_text += symbol.text
                            para_text += " "

                        para_text = para_text.strip()

                        if para_text:
                            sentences.append(
                                {
                                    "text": para_text,
                                    "bounding_box": {
                                        "x": x,
                                        "y": y,
                                        "w": w,
                                        "h": h,
                                    },
                                }
                            )

        # Sort by position (top-to-bottom, left-to-right)
        sentences.sort(
            key=lambda s: (s["bounding_box"]["y"] // 20, s["bounding_box"]["x"])
        )

        # Merge into full sentences based on proximity and punctuation
        merged_sentences = []
        current_sentence = None

        for sent in sentences:
            if current_sentence is None:
                current_sentence = {
                    "text": sent["text"],
                    "bounding_box": sent["bounding_box"],
                }
            else:
                # Check if should merge
                prev_y = current_sentence["bounding_box"]["y"]
                curr_y = sent["bounding_box"]["y"]
                y_diff = abs(curr_y - prev_y)

                # Same line or close (within 30px)
                if y_diff < 30:
                    # Check for punctuation
                    if current_sentence["text"][-1] in ".!?":
                        merged_sentences.append(current_sentence)
                        current_sentence = {
                            "text": sent["text"],
                            "bounding_box": sent["bounding_box"],
                        }
                    else:
                        # Merge
                        current_sentence["text"] += " " + sent["text"]
                        # Update bbox to cover both
                        prev_box = current_sentence["bounding_box"]
                        curr_box = sent["bounding_box"]
                        current_sentence["bounding_box"] = {
                            "x": min(prev_box["x"], curr_box["x"]),
                            "y": min(prev_box["y"], curr_box["y"]),
                            "w": max(
                                prev_box["x"] + prev_box["w"],
                                curr_box["x"] + curr_box["w"],
                            )
                            - min(prev_box["x"], curr_box["x"]),
                            "h": max(
                                prev_box["y"] + prev_box["h"],
                                curr_box["y"] + curr_box["h"],
                            )
                            - min(prev_box["y"], curr_box["y"]),
                        }
                else:
                    merged_sentences.append(current_sentence)
                    current_sentence = {
                        "text": sent["text"],
                        "bounding_box": sent["bounding_box"],
                    }

        if current_sentence:
            merged_sentences.append(current_sentence)

        # Map colors to each sentence based on bounding box
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        if img is not None:
            for sent in merged_sentences:
                bbox = sent["bounding_box"]
                x, y = max(0, bbox["x"]), max(0, bbox["y"])
                w = min(bbox["w"], img.shape[1] - x)
                h = min(bbox["h"], img.shape[0] - y)

                if w > 0 and h > 0:
                    roi = img[y : y + h, x : x + w]

                    # Get dominant color (background)
                    avg_color = roi.mean(axis=(0, 1))
                    bg_b, bg_g, bg_r = (
                        int(avg_color[0]),
                        int(avg_color[1]),
                        int(avg_color[2]),
                    )
                    bg_sum = bg_r + bg_g + bg_b

                    # Get all unique colors
                    unique_rgb = np.array(
                        [
                            list(c[::-1].astype(int))
                            for c in np.unique(roi.reshape(-1, 3), axis=0)
                        ]
                    )

                    def get_text_color_for_region(roi_region):
                        unique = np.unique(roi_region.reshape(-1, 3), axis=0)
                        if len(unique) == 0:
                            return None
                        unique_rgb = np.array(
                            [list(c[::-1].astype(int)) for c in unique]
                        )

                        # Get background (edge pixels)
                        edge_pixels = np.vstack(
                            [
                                roi_region[:5, :].reshape(-1, 3),
                                roi_region[-5:, :].reshape(-1, 3),
                                roi_region[:, :5].reshape(-1, 3),
                                roi_region[:, -5:].reshape(-1, 3),
                            ]
                        )
                        bg_color = np.median(edge_pixels, axis=0)
                        bg_r, bg_g, bg_b = (
                            int(bg_color[2]),
                            int(bg_color[1]),
                            int(bg_color[0]),
                        )
                        bg_sum = bg_r + bg_g + bg_b

                        # Filter out colors too close to white (anti-aliased edges)
                        # and colors too close to background
                        filtered_colors = []
                        for c in unique_rgb:
                            r, g, b = c
                            brightness = r + g + b
                            # Skip if too bright (>600 - likely anti-aliased)
                            if brightness > 600:
                                continue
                            # Skip if too close to background
                            diff = abs(r - bg_r) + abs(g - bg_g) + abs(b - bg_b)
                            if diff < 30:
                                continue
                            filtered_colors.append(c)

                        if len(filtered_colors) == 0:
                            # Fallback: find darkest or brightest from unique
                            dark_sum = unique_rgb.min()
                            light_sum = unique_rgb.max()
                            if abs(light_sum - bg_sum) > abs(dark_sum - bg_sum):
                                return unique_rgb[unique_rgb.argmax()]
                            else:
                                return unique_rgb[unique_rgb.argmin()]

                        filtered = np.array(filtered_colors)

                        # Find color that contrasts most with background
                        contrasts = []
                        for c in filtered:
                            r, g, b = c
                            diff = abs(r - bg_r) + abs(g - bg_g) + abs(b - bg_b)
                            contrasts.append(diff)

                        return filtered[np.argmax(contrasts)]

                    # For wide regions, check left and right separately
                    if w > 300:
                        # Split into left and right halves
                        mid = w // 2
                        roi_left = roi[:, :mid]
                        roi_right = roi[:, mid:]

                        color_left = get_text_color_for_region(roi_left)
                        color_right = get_text_color_for_region(roi_right)

                        # Split text into words and assign colors
                        words = sent["text"].split()
                        if len(words) > 1:
                            # Estimate which words are on left vs right
                            words_per_side = len(words) // 2
                            left_text = " ".join(words[:words_per_side])
                            right_text = " ".join(words[words_per_side:])

                            colors = []
                            if color_left is not None:
                                colors.append(
                                    {
                                        "text": left_text,
                                        "hex": self.rgb_to_hex(
                                            int(color_left[0]),
                                            int(color_left[1]),
                                            int(color_left[2]),
                                        ),
                                        "r": int(color_left[0]),
                                        "g": int(color_left[1]),
                                        "b": int(color_left[2]),
                                    }
                                )
                            if color_right is not None:
                                colors.append(
                                    {
                                        "text": right_text,
                                        "hex": self.rgb_to_hex(
                                            int(color_right[0]),
                                            int(color_right[1]),
                                            int(color_right[2]),
                                        ),
                                        "r": int(color_right[0]),
                                        "g": int(color_right[1]),
                                        "b": int(color_right[2]),
                                    }
                                )
                            sent["text_color"] = colors
                        else:
                            # Single word, use left color
                            if color_left is not None:
                                sent["text_color"] = {
                                    "hex": self.rgb_to_hex(
                                        int(color_left[0]),
                                        int(color_left[1]),
                                        int(color_left[2]),
                                    ),
                                    "r": int(color_left[0]),
                                    "g": int(color_left[1]),
                                    "b": int(color_left[2]),
                                }
                    else:
                        # Single region
                        text_color = get_text_color_for_region(roi)
                        if text_color is not None:
                            sent["text_color"] = {
                                "hex": self.rgb_to_hex(
                                    int(text_color[0]),
                                    int(text_color[1]),
                                    int(text_color[2]),
                                ),
                                "r": int(text_color[0]),
                                "g": int(text_color[1]),
                                "b": int(text_color[2]),
                            }

                    sent["background_color"] = {
                        "hex": self.rgb_to_hex(bg_r, bg_g, bg_b),
                        "r": bg_r,
                        "g": bg_g,
                        "b": bg_b,
                    }

        # Step 6: Return Results
        return {
            "full_text": full_text,
            "sentences": merged_sentences,
            "text_colors": text_colors,
            "background_colors": background_colors,
            "confidence": text_response.text_property_detection_confidence
            if hasattr(text_response, "text_property_detection_confidence")
            else None,
        }

    # def extract_images_from_pptx(
    #     self,
    #     pptx_path: str,
    #     output_folder: str,
    #     user_type: str = "org"
    # ) -> list[str]:
    #     import os, json
    #     from pptx import Presentation
    #     from pptx.enum.shapes import MSO_SHAPE_TYPE

    #     if not os.path.isfile(pptx_path):
    #         raise FileNotFoundError(f"File not found: {pptx_path}")

    #     os.makedirs(output_folder, exist_ok=True)
    #     prs = Presentation(pptx_path)

    #     client = self._build_vision_client(user_type=user_type)

    #     image_paths = []
    #     for s_idx, slide in enumerate(prs.slides, start=1):
    #         for sh_idx, shape in enumerate(slide.shapes, start=1):
    #             if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
    #                 img = shape.image
    #                 ext = img.ext  # 'png' / 'jpeg' etc

    #                 img_path = os.path.join(output_folder, f"slide{s_idx}_img{sh_idx}.{ext}")
    #                 with open(img_path, "wb") as f:
    #                     f.write(img.blob)

    #                 # analysis json next to image
    #                 analysis = self.extract_full_image_analysis(
    #                     image_bytes=img.blob,
    #                     client=client
    #                 )
    #                 json_path = os.path.splitext(img_path)[0] + "_analysis.json"
    #                 with open(json_path, "w", encoding="utf-8") as jf:
    #                     json.dump(analysis, jf, indent=4, ensure_ascii=False)

    #                 image_paths.append(img_path)

    #     return image_paths
